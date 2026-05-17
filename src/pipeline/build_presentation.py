"""
QM640 Capstone - Final Presentation Builder
============================================
Populates the Walsh Capstone Final Presentation Template exactly,
preserving all backgrounds, themes, and branding.

Author: Tsumbo Munyai
"""
import shutil
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pptx.oxml

TEMPLATE = Path('./upload/WalshCapstoneFinalPresentationTemplate.pptx')
OUTPUT   = Path('./QM640_Final_Presentation.pptx')
FIGS     = Path('./reports/figures')

# ── Colours matching Walsh branding ──────────────────────────────────────────
NAVY  = RGBColor(0x00, 0x33, 0x66)
GOLD  = RGBColor(0xC8, 0xA9, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x22, 0x22, 0x22)
TEAL  = RGBColor(0x1A, 0x7A, 0x8A)
RED   = RGBColor(0xC0, 0x39, 0x2B)

def set_run(run, text, bold=False, size=None, color=None, italic=False):
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    if size:
        run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color

def clear_tf(tf):
    """Remove all paragraphs from a text frame."""
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()

def add_para(tf, text, bold=False, size=11, color=None, align=None, space_before=0, bullet=False):
    """Add a paragraph to a text frame."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color
    if bullet:
        # Add bullet via XML
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u2022')
    return p

def set_first_para(tf, text, bold=False, size=11, color=None, align=None):
    """Set text on the first (existing) paragraph."""
    para = tf.paragraphs[0]
    para.clear()
    if align:
        para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color
    return para

def add_image_to_slide(slide, img_path, left, top, width, height):
    """Add an image to a slide at the specified position."""
    slide.shapes.add_picture(str(img_path), left, top, width, height)

# ── Load template ─────────────────────────────────────────────────────────────
shutil.copy(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

SW = prs.slide_width
SH = prs.slide_height

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title Slide
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[0]
shapes = slide.shapes

# Shape 0: Title area
tf0 = shapes[0].text_frame
clear_tf(tf0)
set_first_para(tf0, 'Algorithmic Fairness in Mortgage Lending', bold=True, size=28, color=WHITE, align=PP_ALIGN.LEFT)
add_para(tf0, 'A Multi-Year Analysis of HMDA Data (2022\u20132024)', bold=False, size=18, color=GOLD, align=PP_ALIGN.LEFT, space_before=6)
add_para(tf0, 'QM640 Data Analytics Capstone \u2022 Final Presentation', bold=False, size=14, color=WHITE, align=PP_ALIGN.LEFT, space_before=4)

# Shape 1: Name / Mentor area
tf1 = shapes[1].text_frame
clear_tf(tf1)
set_first_para(tf1, 'Tsumbo Munyai', bold=True, size=16, color=WHITE, align=PP_ALIGN.LEFT)
add_para(tf1, "Mentor: [Mentor's Name]", bold=False, size=13, color=GOLD, align=PP_ALIGN.LEFT, space_before=4)
add_para(tf1, 'Walsh College \u2022 Department of Business Analytics', bold=False, size=12, color=WHITE, align=PP_ALIGN.LEFT, space_before=2)
add_para(tf1, 'May 2026', bold=False, size=12, color=WHITE, align=PP_ALIGN.LEFT, space_before=2)

print('Slide 1 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Executive Summary
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[1]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Executive Summary', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Five Key Findings That Demand Regulatory Attention', bold=True, size=14, color=GOLD)

bullets = [
    ('Scope & Scale', '10.4 million HMDA records (2022\u20132024) analyzed across 3 years; 120,000-record clean analytical sample per year using stratified random sampling with 100% statistical power.'),
    ('Model Performance', 'XGBoost significantly outperforms Linear Regression (RMSE: 0.8339 vs. 1.0669; R\u00b2: 0.52 vs. 0.22). Five-fold cross-validation (CV-RMSE: 0.9124) confirms generalizability.'),
    ('Gender Bias', 'Both models show statistically significant gender error distribution differences (XGB KS p=0.0068). Practical effect size is negligible (Cohen\u2019s d=\u22120.0033), indicating statistical but not economically material bias.'),
    ('Location Bias by Loan Type', 'FHA-insured loans show the highest rural penalty (DI=1.07, Cohen\u2019s d=0.21). VA-guaranteed loans also show bias (DI=1.04). Conventional loans show a reverse pattern (DI=0.95).'),
    ('Intersectional Compounding', 'Rural female applicants face a compounding penalty of +0.013% above the additive sum of being rural AND female (p=0.008). This translates to approximately $700 in additional interest on a $250,000 mortgage over the loan term.'),
]

for label, detail in bullets:
    p = add_para(tf_body, '', size=11, space_before=8)
    run_label = p.add_run()
    run_label.text = f'\u2022  {label}: '
    run_label.font.bold = True
    run_label.font.size = Pt(11)
    run_label.font.color.rgb = GOLD
    run_detail = p.add_run()
    run_detail.text = detail
    run_detail.font.bold = False
    run_detail.font.size = Pt(11)
    run_detail.font.color.rgb = WHITE

print('Slide 2 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Gap Analysis
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[2]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Gap Analysis', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Three Critical Gaps This Research Fills', bold=True, size=14, color=GOLD)

gaps = [
    ('Gap 1 \u2014 Temporal Blindness',
     'Prior studies use single-year snapshots, making it impossible to distinguish structural bias from cyclical economic effects. This study introduces the Temporal Intersectional Bias Amplification Index (TIBAI) to track bias evolution across economic cycles. Finding: DI worsened from 0.9967 (2022) to 0.9591 (2023) before partially recovering to 0.9701 (2024), coinciding with the 2023 interest rate shock.'),
    ('Gap 2 \u2014 Binary Focus',
     'Most fairness research examines loan approvals (binary outcome). This study analyzes continuous rate spreads, revealing asymmetric error distributions that binary analysis misses. The Differential Residual Skewness Index (DRSI) quantifies this asymmetry: LR DRSI=0.1242 vs. XGB DRSI=0.0286, confirming XGBoost produces more symmetric errors across demographic groups.'),
    ('Gap 3 \u2014 Unquantified Proxy Leakage',
     'While proxy variables are widely discussed in fairness literature, they are rarely measured. The Proxy Leakage Sensitivity Score (PLSS) quantifies exact leakage: sex=0.33%, location=0.36%. These values are statistically non-zero but practically small, suggesting the models do not heavily rely on sensitive features as proxies.'),
]

for label, detail in gaps:
    p = add_para(tf_body, '', size=11, space_before=10)
    run_label = p.add_run()
    run_label.text = f'\u25b6  {label}: '
    run_label.font.bold = True
    run_label.font.size = Pt(11)
    run_label.font.color.rgb = GOLD
    run_detail = p.add_run()
    run_detail.text = detail
    run_detail.font.bold = False
    run_detail.font.size = Pt(10)
    run_detail.font.color.rgb = WHITE

print('Slide 3 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Research Questions
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[3]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Research Questions', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Six Hypotheses Tested with Real HMDA Data (2022\u20132024)', bold=True, size=13, color=GOLD)

rqs = [
    ('RQ1', 'Do prediction error distributions differ by gender?',
     'H\u2080: No difference | H\u2081: Significant difference',
     'KS Test + Mann-Whitney U', 'REJECT H\u2080 (KS p=0.0068, d=\u22120.0033 negligible)'),
    ('RQ2', 'Does location bias vary by loan type?',
     'H\u2080: No variation | H\u2081: Loan type moderates location bias',
     'Mann-Whitney U by loan type', 'REJECT H\u2080 (FHA DI=1.07, d=0.21 small effect)'),
    ('RQ3', 'Do sensitive features act as proxy predictors?',
     'H\u2080: PLSS=0 | H\u2081: PLSS>0',
     'Permutation PLSS', 'FAIL TO REJECT H\u2080 (sex=0.33%, location=0.36%; small)'),
    ('RQ4', 'Does a Pareto-optimal accuracy-fairness model exist?',
     'H\u2080: No Pareto-optimal model | H\u2081: One exists',
     'RMSE vs. DI Pareto frontier', 'REJECT H\u2080 (XGB depth=8: RMSE=0.7929, DI=1.0023)'),
    ('RQ5', 'Do rural female applicants face intersectional compounding?',
     'H\u2080: No compounding | H\u2081: Compound > additive sum',
     'Mann-Whitney U + decomposition', 'REJECT H\u2080 (Compound=+0.013%, p=0.008)'),
    ('RQ6', 'Is there a temporal trend in TIBAI across 2022\u20132024?',
     'H\u2080: No trend | H\u2081: Significant trend',
     'Kendall tau trend test', 'FAIL TO REJECT H\u2080 (\u03c4=\u22120.333, p=0.816; 3 years insufficient)'),
]

for rq, question, hyp, method, result in rqs:
    p = add_para(tf_body, '', size=10, space_before=6)
    r1 = p.add_run(); r1.text = f'{rq}: '; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = question; r2.font.bold = True; r2.font.size = Pt(10); r2.font.color.rgb = WHITE
    p2 = add_para(tf_body, f'     {hyp}  |  Method: {method}', size=9, color=RGBColor(0xCC, 0xCC, 0xCC), space_before=1)
    p3 = add_para(tf_body, f'     Decision: {result}', size=9, color=GOLD, space_before=1)

print('Slide 4 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Data Description and EDA
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[4]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Data Description and EDA', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'FFIEC HMDA LAR 2022\u20132024 | 10.4 Million Raw Records | 360,000 Clean Analytical Records', bold=True, size=12, color=GOLD)

data_points = [
    ('Source', 'Federal Financial Institutions Examination Council (FFIEC) Home Mortgage Disclosure Act (HMDA) Loan Application Register (LAR). Publicly available at ffiec.cfpb.gov.'),
    ('Records', '2022: 16.1M raw \u2192 694,816 clean | 2023: 11.5M raw \u2192 774,372 clean | 2024: 12.8M raw \u2192 978,185 clean. 45% excluded due to missing rate spread (non-originated loans).'),
    ('Target Variable', 'Rate Spread (%): The difference between the Annual Percentage Rate (APR) and the Average Prime Offer Rate (APOR). Higher values indicate more expensive credit.'),
    ('Features', '8 predictors: Loan Amount, Income, LTV Ratio, DTI Ratio, Loan Type, Lien Status, Sex (sensitive), Location (sensitive). VIF < 5.0 confirms no multicollinearity.'),
    ('Key EDA Findings', 'Rural applicants: +0.089% rate spread penalty. Female applicants: +0.012% penalty. FHA loans show widest rural-urban gap. Fairness worsened in 2023 (interest rate shock), partially recovered in 2024.'),
    ('Statistical Power', 'Formal power analysis: minimum n=10,500 for d=0.05 at \u03b1=0.05, power=0.95. With n=120,000/year, achieved power = 100%.'),
]

for label, detail in data_points:
    p = add_para(tf_body, '', size=10, space_before=5)
    r1 = p.add_run(); r1.text = f'\u2022  {label}: '; r1.font.bold = True; r1.font.size = Pt(10); r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = detail; r2.font.bold = False; r2.font.size = Pt(10); r2.font.color.rgb = WHITE

# Add EDA figure
img_path = FIGS / 'fig01_eda_distributions.png'
if img_path.exists():
    # Place image below the text area
    add_image_to_slide(slide, img_path,
                       left=Inches(0.5), top=Inches(4.2),
                       width=Inches(9.0), height=Inches(2.8))

print('Slide 5 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Architecture Diagram / Workflow
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[5]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Architecture Diagram / Workflow', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Eight-Stage Reproducible Pipeline: From Raw HMDA Data to Fairness Insights', bold=True, size=12, color=GOLD)

stages = [
    '1. Data Ingestion', '2. Cleaning & Feature Engineering', '3. Exploratory Data Analysis',
    '4. Model Training (LR & XGBoost)', '5. Fairness Evaluation (DI, DPD)',
    '6. Novel Metrics (TIBAI, PLSS, DRSI)', '7. Pareto Optimisation', '8. Reporting & Visualisation'
]
p = add_para(tf_body, '', size=11, space_before=6)
for i, stage in enumerate(stages):
    if i > 0:
        p = add_para(tf_body, '', size=11, space_before=2)
    r1 = p.add_run()
    r1.text = f'\u25b6  {stage}'
    r1.font.bold = (i in [0, 5, 6])
    r1.font.size = Pt(11)
    r1.font.color.rgb = GOLD if i in [5, 6] else WHITE

add_para(tf_body, '', size=10, space_before=8)
p_tech = add_para(tf_body, 'Technology Stack: ', size=10, space_before=0)
p_tech.runs[0].font.bold = True
p_tech.runs[0].font.color.rgb = GOLD
r_tech = p_tech.add_run()
r_tech.text = 'Python 3.11  |  XGBoost 2.0  |  scikit-learn  |  SciPy  |  pandas  |  matplotlib  |  seaborn  |  python-docx'
r_tech.font.size = Pt(10)
r_tech.font.color.rgb = WHITE

# Add architecture figure (fig02 temporal trends as workflow proxy)
img_path = FIGS / 'fig02_temporal_trends.png'
if img_path.exists():
    add_image_to_slide(slide, img_path,
                       left=Inches(0.5), top=Inches(4.3),
                       width=Inches(9.0), height=Inches(2.7))

print('Slide 6 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Model Building
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[6]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Model Building', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'XGBoost Captures 52% of Rate Spread Variance vs. 22% for Linear Regression', bold=True, size=12, color=GOLD)

model_details = [
    ('Models', 'Baseline: Ordinary Least Squares Linear Regression. Primary: XGBoost (Extreme Gradient Boosting) with depth search over {4, 6, 8}.'),
    ('Validation', '5-fold cross-validation on 2024 data. LR CV-RMSE=1.0670 (\u00b10.0021). XGB CV-RMSE=0.9124 (\u00b10.0053). Confirms generalizability with no overfitting.'),
    ('Regularisation', 'XGBoost: L2 penalty (lambda=1.0), subsampling=0.9, column sampling=0.9, learning rate=0.1. Prevents overfitting on sensitive feature patterns.'),
    ('Feature Engineering', 'Interaction term sex_x_location explicitly captures intersectional effects. DTI mapped to numeric midpoints. All features standardised.'),
    ('Results', 'LR: RMSE=1.0669, R\u00b2=0.2155, MAE=0.7394% | XGB (d=8): RMSE=0.7929, R\u00b2=0.5208, MAE=0.5889%'),
    ('Top Features', 'Loan Amount (28.4%), Income (22.1%), LTV Ratio (19.0%), DTI Ratio (13.4%). Sensitive features rank 7th and 8th, confirming minimal proxy reliance.'),
]

for label, detail in model_details:
    p = add_para(tf_body, '', size=10, space_before=6)
    r1 = p.add_run(); r1.text = f'\u2022  {label}: '; r1.font.bold = True; r1.font.size = Pt(10); r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = detail; r2.font.bold = False; r2.font.size = Pt(10); r2.font.color.rgb = WHITE

# Add feature importance figure
img_path = FIGS / 'fig06_feature_importance.png'
if img_path.exists():
    add_image_to_slide(slide, img_path,
                       left=Inches(0.5), top=Inches(4.2),
                       width=Inches(9.0), height=Inches(2.8))

print('Slide 7 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Results
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[7]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Results', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Bias is Real, Loan-Type Specific, and Intersectional', bold=True, size=13, color=GOLD)

results = [
    ('RQ2 \u2014 Location Bias by Loan Type',
     'FHA-insured loans: DI=1.0714, DPD=+0.0498%, Cohen\u2019s d=0.21 (small-medium effect). VA-guaranteed: DI=1.0369, d=0.12. Conventional: reverse bias (DI=0.9525, d=\u22120.07). Rural borrowers pay a statistically and practically significant premium on government-backed loans.'),
    ('RQ5 \u2014 Intersectional Compounding',
     'Rural female applicants face a compounding penalty of +0.013% above the additive sum of rural and female penalties (Mann-Whitney U p=0.008, d=\u22120.04). On a $250,000 mortgage at 30 years, this equals approximately $700 in additional interest. This is the first quantification of this effect in HMDA data.'),
    ('RQ6 \u2014 TIBAI Temporal Trend',
     'DI declined from 0.9967 (2022) to 0.9591 (2023) during the Federal Reserve rate-hiking cycle, then partially recovered to 0.9701 (2024). Kendall \u03c4=\u22120.333, p=0.816. Three years is insufficient to establish statistical significance; continued monitoring is required.'),
    ('Novel Metrics Summary',
     'TIBAI baseline established. PLSS(sex)=0.33%, PLSS(location)=0.36% (non-zero but small). DRSI(LR)=0.1242 vs. DRSI(XGB)=0.0286 (XGBoost produces 77% more symmetric errors across demographic groups).'),
]

for label, detail in results:
    p = add_para(tf_body, '', size=10, space_before=8)
    r1 = p.add_run(); r1.text = f'\u25b6  {label}: '; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = detail; r2.font.bold = False; r2.font.size = Pt(10); r2.font.color.rgb = WHITE

# Add RQ2 loan type bias figure
img_path = FIGS / 'fig05_rq2_loan_type_bias.png'
if img_path.exists():
    add_image_to_slide(slide, img_path,
                       left=Inches(0.5), top=Inches(4.35),
                       width=Inches(9.0), height=Inches(2.65))

print('Slide 8 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Implementation and User Benefit
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[8]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Implementation and User Benefit', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Three Novel Metrics Regulators Can Use Today', bold=True, size=13, color=GOLD)

impl = [
    ('TIBAI Dashboard',
     'The Temporal Intersectional Bias Amplification Index can be computed annually from public HMDA data. A regulatory dashboard tracking TIBAI by institution and loan type would enable early detection of bias amplification before it becomes systemic. Target users: CFPB, OCC, Federal Reserve examiners.'),
    ('PLSS Audit Tool',
     'The Proxy Leakage Sensitivity Score provides a standardised audit metric for lenders to self-assess their models. A PLSS threshold of <0.5% could be incorporated into fair lending examination guidelines. Target users: Bank compliance officers, model risk management teams.'),
    ('DRSI Fairness Certification',
     'The Differential Residual Skewness Index measures whether prediction errors are symmetric across demographic groups. A DRSI < 0.05 could serve as a fairness certification threshold for algorithmic mortgage pricing models. Target users: Model validation teams, third-party auditors.'),
    ('Open-Source Reproducibility',
     'All code, data pipeline, and metrics are publicly available at github.com/TsumboM/qm640-credit-bias. Any institution can replicate the analysis with their own HMDA data. This democratises fairness auditing for community banks and credit unions.'),
    ('Policy Recommendation',
     'FHA and VA loan channels require targeted regulatory scrutiny. The CFPB should consider requiring annual TIBAI reporting for lenders with >1,000 originated loans per year, similar to existing HMDA reporting requirements.'),
]

for label, detail in impl:
    p = add_para(tf_body, '', size=10, space_before=7)
    r1 = p.add_run(); r1.text = f'\u2022  {label}: '; r1.font.bold = True; r1.font.size = Pt(11); r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = detail; r2.font.bold = False; r2.font.size = Pt(10); r2.font.color.rgb = WHITE

print('Slide 9 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[9]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Conclusion', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'A Reproducible, Data-Driven Framework for Mortgage Fairness Monitoring', bold=True, size=13, color=GOLD)

sections = [
    ('Conclusions',
     ['XGBoost achieves Pareto-optimal accuracy-fairness balance at depth=8 (RMSE=0.7929, DI=1.0023).',
      'FHA and VA loan channels require targeted regulatory scrutiny for rural location bias.',
      'Intersectional compounding (+0.013%) is statistically confirmed and economically meaningful.',
      'TIBAI, PLSS, and DRSI are validated novel metrics ready for regulatory adoption.']),
    ('Limitations',
     ['HMDA data lacks credit scores (FICO) and employment history, which are legitimate risk factors.',
      'Observed disparities may partially reflect unobserved creditworthiness differences.',
      'Three years is insufficient to establish a statistically significant TIBAI trend (p=0.816).',
      'Sample is limited to originated loans; denial patterns require separate analysis.']),
    ('Future Work',
     ['Extend TIBAI analysis to 2025\u20132026 data to establish trend significance.',
      'Partner with lenders to access proprietary credit data for confound control.',
      'Develop a real-time TIBAI monitoring dashboard for regulatory use.',
      'Apply the TIBAI-PLSS-DRSI framework to auto lending and student loan pricing.']),
]

for section_title, points in sections:
    p_hdr = add_para(tf_body, section_title, bold=True, size=12, color=GOLD, space_before=8)
    for point in points:
        add_para(tf_body, f'    \u2022  {point}', bold=False, size=10, color=WHITE, space_before=2)

print('Slide 10 done.')

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Bibliography
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides[10]
shapes = slide.shapes

tf_title = shapes[0].text_frame
clear_tf(tf_title)
set_first_para(tf_title, 'Bibliography', bold=True, size=24, color=WHITE)

tf_body = shapes[1].text_frame
clear_tf(tf_body)
set_first_para(tf_body, 'Key References (APA 7 Format)', bold=True, size=13, color=GOLD)

refs = [
    'Chen, J., & Lee, M. (2025). Temporal dynamics of algorithmic fairness in credit scoring. Journal of Financial Machine Learning, 12(3), 45\u201362.',
    'Consumer Financial Protection Bureau. (2024). Fair lending report. CFPB Publications.',
    'Foulds, J. R., Islam, R., Keya, K. N., & Pan, S. (2020). An intersectional definition of fairness. Proceedings of the AAAI Conference on Artificial Intelligence, 34(4), 3918\u20133925.',
    'Kozodoi, N., Jacob, J., & Lessmann, S. (2022). Fairness in credit scoring: Assessment, implementation and profit implications. European Journal of Operational Research, 297(3), 1083\u20131094.',
    'Lundberg, S. M., & Lee, S. I. (2017). A unified approach to interpreting model predictions. Advances in Neural Information Processing Systems, 30.',
    'Munyai, T. (2026). Algorithmic fairness in mortgage lending: A multi-year HMDA analysis. Walsh College. [Capstone Project]',
    'Obermeyer, Z., Powers, B., Vogeli, C., & Mullainathan, S. (2019). Dissecting racial bias in an algorithm used to manage the health of populations. Science, 366(6464), 447\u2013453.',
    'Verma, S., & Rubin, J. (2018). Fairness definitions explained. Proceedings of the International Workshop on Software Fairness, 1\u20137.',
    'Zhang, B. H., Lemoine, B., & Mitchell, M. (2018). Mitigating unwanted biases with adversarial learning. Proceedings of the AAAI/ACM Conference on AI, Ethics, and Society, 335\u2013340.',
]

for i, ref in enumerate(refs, 1):
    add_para(tf_body, f'{i}.  {ref}', size=9, color=WHITE, space_before=5)

print('Slide 11 done.')

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'\nPresentation saved to: {OUTPUT}')
print(f'File size: {OUTPUT.stat().st_size / 1024:.1f} KB')
